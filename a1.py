import pytesseract
from PIL import Image
from pptx import Presentation
import os

# Set path to Tesseract (Make sure to update with your installation path)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # For Windows, update if necessary

# Function to extract text from PNG image
def extract_text_from_image(image_path):
    img = Image.open(image_path)
    return pytesseract.image_to_string(img)

# Function to create a PowerPoint presentation from the extracted text
def create_ppt_from_text(text):
    prs = Presentation()
    
    # Add title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Extracted Text"

    # Add slide with content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    content = slide.shapes.placeholders[1]
    content.text = text

    # Save the PowerPoint presentation
    prs.save('output.pptx')

# Function to process all PNG files in multiple folders
def process_images_in_multiple_folders(parent_directory):
    # Loop through all subdirectories (dataset-001, dataset-002, etc.)
    for folder_name in os.listdir(parent_directory):
        folder_path = os.path.join(parent_directory, folder_name)
        
        if os.path.isdir(folder_path):  # Check if it's a directory
            print(f"Processing folder: {folder_name}")
            
            # Loop through all files in the folder
            for filename in os.listdir(folder_path):
                if filename.endswith(".png"):  # Process only PNG files
                    image_path = os.path.join(folder_path, filename)
                    text = extract_text_from_image(image_path)  # Extract text
                    create_ppt_from_text(text)  # Create PowerPoint slide

# Replace with the parent directory where your dataset folders are stored
parent_directory = r'C:\SmartBoardProject\dataset'  # Update with the correct parent folder path

# Process all images in all folders
process_images_in_multiple_folders(parent_directory)