from pptx import Presentation
from transformers import TrOCRProcessor, VisionEncoderDecoderModel
from PIL import Image
from ultralytics import YOLO
import openai
import torch
import cv2

# Set API Key
openai.api_key = 'AIzaSyCBP0gFI0gnfFkg_JhDrWkpJx5LUro838Q'

# Load models
processor = TrOCRProcessor.from_pretrained('microsoft/trocr-base-handwritten')
trocr_model = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-base-handwritten')
yolo_model = YOLO('yolov8n.pt')  # Use your trained model if available

def preprocess_image(image_path):
    img = cv2.imread(image_path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    output_path = 'preprocessed.jpg'
    cv2.imwrite(output_path, thresh)
    return output_path

def extract_text(image_path):
    image = Image.open(image_path).convert("RGB")
    pixel_values = processor(images=image, return_tensors="pt").pixel_values
    with torch.no_grad():
        generated_ids = trocr_model.generate(pixel_values)
    return processor.batch_decode(generated_ids, skip_special_tokens=True)[0]

def detect_diagrams(image_path):
    results = yolo_model.predict(source=image_path, save=False)
    return "Detected diagrams: " + str([result.boxes for result in results])

def generate_slide_structure(text, diagrams):
    prompt = f"""
    Create a structured PowerPoint outline from the following:

    Text: {text}
    Diagrams: {diagrams}
    """

    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )

    return response['choices'][0]['message']['content']

def create_ppt(slide_data, output_file='final_slides.pptx'):
    prs = Presentation()
    for slide in slide_data:
        slide_layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(slide_layout)
        s.shapes.title.text = slide['title']
        s.placeholders[1].text = '\n'.join(slide['points'])
    prs.save(output_file)

def main():
    raw_image = 'sample_whiteboard.jpg'
    
    # 1. Preprocess
    preprocessed = preprocess_image(raw_image)

    # 2. Extract text
    text = extract_text(preprocessed)
    
    # 3. Detect diagrams
    diagrams = detect_diagrams(raw_image)
    
    # 4. Generate slide outline
    structured_text = generate_slide_structure(text, diagrams)
    
    # Convert structured text to list of slides manually for this demo
    slides = [
        {'title': 'Introduction', 'points': ['Overview of AI', 'Importance of AI']},
        {'title': 'Applications', 'points': ['Healthcare', 'Finance', 'Education']}
    ]

    # 5. Create PPT
    create_ppt(slides)

    print("All Done!")

if __name__ == '__main__':
    main()
