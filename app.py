from flask import Flask, render_template, request, send_file, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import cv2
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import uuid
from werkzeug.utils import secure_filename
import logging
from concurrent.futures import ThreadPoolExecutor

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

app = Flask(__name__)
UPLOAD_FOLDER = os.path.normpath('uploads')
OUTPUT_FOLDER = os.path.normpath('C:/DeckifyOutput')  # Moved outside OneDrive to avoid sync issues
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # Limit uploads to 16MB

# Create directories if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Set Tesseract path dynamically
try:
    pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_CMD', 'tesseract')
    logging.info(f"Tesseract version: {pytesseract.get_tesseract_version()}")
except Exception as e:
    logging.error(f"Tesseract not found or misconfigured: {e}")
    raise RuntimeError("Tesseract OCR is not installed or not configured properly.")

# Allowed file extensions for images
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'bmp', 'gif'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Preprocessing function to improve OCR accuracy
def preprocess_image(img_path):
    logging.info(f"Preprocessing image: {img_path}")
    try:
        img = Image.open(img_path)
        logging.info(f"Image opened successfully: {img_path}")
        
        # Resize image if too large
        max_size = (1000, 1000)  # Max width/height
        img.thumbnail(max_size, Image.Resampling.LANCZOS)
        
        # Convert to grayscale
        img = img.convert('L')
        
        # Reduce noise with a median filter
        img = img.filter(ImageFilter.MedianFilter(size=3))
        
        # Sharpen the image
        img = img.filter(ImageFilter.SHARPEN)
        
        # Increase contrast
        img = ImageEnhance.Contrast(img).enhance(3)
        
        # Adjust brightness
        img = ImageEnhance.Brightness(img).enhance(1.5)
        
        # Apply thresholding to binarize the image (black and white)
        img = img.point(lambda p: 255 if p > 128 else 0)
        
        return img
    except Exception as e:
        logging.error(f"Error preprocessing image {img_path}: {e}")
        raise

# Function to extract text from an image
def extract_text(img_path):
    logging.info(f"Extracting text from: {img_path}")
    try:
        img = preprocess_image(img_path)
        text = pytesseract.image_to_string(img, config='--psm 6')
        logging.info(f"Text extracted from {img_path}: {text[:100]}...")  # Log first 100 chars
        return text.strip() if text.strip() else "No text detected in this image."
    except Exception as e:
        logging.error(f"Error extracting text from {img_path}: {e}")
        return f"Error: Unable to extract text from this image."

# Function to add a footer to a slide
def add_footer(slide, footer_text):
    try:
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        text_frame = footer.text_frame
        p = text_frame.add_paragraph()
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)  # Gray color for footer
    except Exception as e:
        logging.error(f"Error adding footer to slide: {e}")

# Function to create a PPT
def create_ppt(image_paths):
    logging.info("Starting PPT creation...")
    try:
        ppt = Presentation()

        # Add Title Slide with Project Name and Logo
        title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = "Smart Board to Slide Deck Converter"
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Converts handwritten board content into presentation slides"

        # Add Logo to the title slide
        logo_path = os.path.join('static', 'logo.png')
        if os.path.exists(logo_path):
            title_slide.shapes.add_picture(logo_path, Inches(5), Inches(3), width=Inches(3), height=Inches(3))
        else:
            logging.warning("Logo file not found at 'static/logo.png'.")

        # Process images in parallel
        with ThreadPoolExecutor() as executor:
            text_results = list(executor.map(extract_text, image_paths))

        # Add content slides with extracted text and logo
        for img_path, extracted_text in zip(image_paths, text_results):
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use layout 1 for safety

            # Display the original image
            try:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.0), width=Inches(4.0), height=Inches(3.0))
            except Exception as e:
                logging.error(f"Error adding image {img_path} to slide: {e}")
                continue

            # Add extracted text to slide
            textbox = slide.shapes.add_textbox(Inches(4.5), Inches(1.0), Inches(4.0), Inches(3.0))
            text_frame = textbox.text_frame
            text_frame.text = extracted_text
            text_frame.word_wrap = True
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(14)

            # Add Logo and Footer
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(8.0), Inches(5.5), width=Inches(1.0), height=Inches(1.0))

            add_footer(slide, "Smart Board to Slide Deck Converter")

        # Generate a unique output filename
        output_filename = f"presentation_{uuid.uuid4()}.pptx"
        output_file = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
        
        # Ensure the output directory is writable
        output_dir = os.path.dirname(output_file)
        os.makedirs(output_dir, exist_ok=True)
        
        # Save the PPT file
        ppt.save(output_file)
        logging.info(f"PPT saved successfully: {output_file}")
        
        # Verify the file exists after saving
        if not os.path.exists(output_file):
            raise FileNotFoundError(f"Failed to save PPT file: {output_file}")
        
        return output_file, output_filename

    except Exception as e:
        logging.error(f"Error creating PPT: {e}")
        raise

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    logging.info("Received upload request.")
    try:
        # Check if files were uploaded
        if 'files' not in request.files:
            logging.error("No files part in the request.")
            return {"error": "No files uploaded."}, 400

        files = request.files.getlist('files')
        if not files or all(file.filename == '' for file in files):
            logging.error("No files selected for upload.")
            return {"error": "No files selected."}, 400

        image_paths = []
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                filename = str(uuid.uuid4()) + os.path.splitext(filename)[-1]
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                logging.info(f"File saved: {filepath}")
                image_paths.append(filepath)
            else:
                logging.warning(f"Invalid file type: {file.filename}")
                return {"error": f"Invalid file type: {file.filename}. Only images are allowed."}, 400

        if not image_paths:
            logging.error("No valid images uploaded.")
            return {"error": "No valid images uploaded."}, 400

        # Create PPT
        output_file, output_filename = create_ppt(image_paths)

        # Verify the PPT file exists before sending
        if not os.path.exists(output_file):
            logging.error(f"PPT file not found: {output_file}")
            raise FileNotFoundError(f"PPT file not found: {output_file}")

        # Clean up uploaded images
        for img_path in image_paths:
            try:
                os.remove(img_path)
                logging.info(f"Deleted uploaded image: {img_path}")
            except Exception as e:
                logging.warning(f"Failed to delete uploaded image {img_path}: {e}")

        logging.info(f"Sending file: {output_filename}")
        return send_file(output_file, as_attachment=True, download_name=output_filename)

    except Exception as e:
        logging.error(f"Error in upload_file: {e}")
        return {"error": f"Failed to process the image: {str(e)}"}, 500

# Test route to verify PPT creation
@app.route('/test_ppt', methods=['GET'])
def test_ppt():
    try:
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide"
        output_file = os.path.join(app.config['OUTPUT_FOLDER'], 'test.pptx')
        ppt.save(output_file)
        logging.info(f"Test PPT saved: {output_file}")
        return send_file(output_file, as_attachment=True, download_name='test.pptx')
    except Exception as e:
        logging.error(f"Error in test_ppt: {e}")
        return {"error": str(e)}, 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)  # Debug enabled for troubleshooting