from pptx import Presentation

def create_ppt(slide_content_list, output_file='smartboard_presentation.pptx'):
    prs = Presentation()

    for slide in slide_content_list:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        s = prs.slides.add_slide(slide_layout)

        title = s.shapes.title
        content = s.placeholders[1]

        title.text = slide['title']
        content.text = '\n'.join(slide['points'])

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example Test
slide_data = [
    {'title': 'Introduction to AI', 'points': ['Definition of AI', 'History of AI']},
    {'title': 'Applications', 'points': ['Healthcare', 'Finance', 'Education']}
]

create_ppt(slide_data)
