import pytesseract
from PIL import Image
import os
from pptx import Presentation

# Path to Tesseract executable (make sure this is correct on your system)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Update the path if necessary

# Function to extract text from an image
def extract_text_from_image(image_path):
    try:
        img = Image.open(image_path)
        custom_config = r'--oem 3 --psm 6'  # OCR engine options for better accuracy
        text = pytesseract.image_to_string(img, config=custom_config)
        return text.strip()
    except Exception as e:
        print(f"Error processing {image_path}: {e}")
        return ""  # Return empty string if error occurs

# Function to create PowerPoint from extracted text
def create_ppt_from_text(text, output_ppt):
    prs = Presentation()
    
    # Add title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Extracted Text"

    # Add content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    content = slide.shapes.placeholders[1]
    content.text = text

    # Save the presentation
    prs.save(output_ppt)
    print(f"PowerPoint saved: {output_ppt}")

# Function to process all images in the dataset folder and create PowerPoint
def process_images_in_folder(dataset_folder, output_folder):
    # Check if the output folder exists, if not, create it
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Iterate through each folder in the dataset
    for folder_name in os.listdir(dataset_folder):
        folder_path = os.path.join(dataset_folder, folder_name)
        if os.path.isdir(folder_path):  # Check if it's a folder
            print(f"Processing folder: {folder_name}")
            
            ppt_filename = os.path.join(output_folder, f"{folder_name}_output.pptx")
            ppt_text = ""

            # Iterate through each PNG file in the folder
            for filename in os.listdir(folder_path):
                if filename.lower().endswith(".png"):
                    image_path = os.path.join(folder_path, filename)
                    print(f"Processing image: {image_path}")
                    text = extract_text_from_image(image_path)  # Extract text
                    ppt_text += text + "\n\n"  # Add the extracted text to the PowerPoint content

            # Create a PowerPoint presentation with the extracted text
            create_ppt_from_text(ppt_text, ppt_filename)

# Main function
if __name__ == "__main__":
    dataset_folder = r'C:\SmartBoardProject\dataset'  # Update this with the path to your dataset folder
    output_folder = r'C:\SmartBoardProject\output'  # Path to save the output PowerPoint files
    
    # Start processing the images
    process_images_in_folder(dataset_folder, output_folder)